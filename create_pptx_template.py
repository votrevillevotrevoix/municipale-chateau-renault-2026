#!/usr/bin/env python3
"""
Créateur de template PowerPoint pour Votre Ville, Votre Voix
Château-Renault 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Charte graphique du site
VERT_FONCE = RGBColor(44, 95, 93)      # #2c5f5d
VERT_CLAIR = RGBColor(127, 209, 174)   # #7fd1ae
VERT_MOYEN = RGBColor(61, 139, 133)    # #3d8b85
BLANC = RGBColor(255, 255, 255)
NOIR = RGBColor(0, 0, 0)

def create_presentation():
    """Crée une présentation PowerPoint avec la charte graphique"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Page de titre avec gradient
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Fond vert foncé
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = VERT_FONCE

    # Titre principal
    title_box = slide1.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "VOTRE VILLE, VOTRE VOIX"
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    title_run = title_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(54)
    title_run.font.bold = True
    title_run.font.color.rgb = VERT_CLAIR
    title_run.font.name = 'Segoe UI'

    # Sous-titre
    subtitle_box = slide1.shapes.add_textbox(
        Inches(0.5), Inches(4.2), Inches(9), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Château-Renault 2026"
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle_run = subtitle_frame.paragraphs[0].runs[0]
    subtitle_run.font.size = Pt(32)
    subtitle_run.font.color.rgb = BLANC
    subtitle_run.font.name = 'Segoe UI'

    # Footer
    footer_box = slide1.shapes.add_textbox(
        Inches(0.5), Inches(6.8), Inches(9), Inches(0.5)
    )
    footer_frame = footer_box.text_frame
    footer_frame.text = "Municipales 15 et 22 mars 2026"
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    footer_run = footer_frame.paragraphs[0].runs[0]
    footer_run.font.size = Pt(18)
    footer_run.font.color.rgb = VERT_CLAIR
    footer_run.font.name = 'Segoe UI'


    # Slide 2: Titre et contenu
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # Fond blanc
    background2 = slide2.background
    fill2 = background2.fill
    fill2.solid()
    fill2.fore_color.rgb = BLANC

    # Bandeau supérieur vert
    header_shape = slide2.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    header_fill = header_shape.fill
    header_fill.solid()
    header_fill.fore_color.rgb = VERT_FONCE
    header_shape.line.fill.background()

    # Titre dans le bandeau
    header_text = header_shape.text_frame
    header_text.text = "Titre de votre diapositive"
    header_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    header_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    header_run = header_text.paragraphs[0].runs[0]
    header_run.font.size = Pt(36)
    header_run.font.bold = True
    header_run.font.color.rgb = BLANC
    header_run.font.name = 'Segoe UI'

    # Zone de contenu
    content_box = slide2.shapes.add_textbox(
        Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Premier paragraphe
    p1 = content_frame.paragraphs[0]
    p1.text = "• Point important 1"
    p1.level = 0
    p1_run = p1.runs[0]
    p1_run.font.size = Pt(24)
    p1_run.font.color.rgb = VERT_FONCE
    p1_run.font.name = 'Segoe UI'

    # Deuxième paragraphe
    p2 = content_frame.add_paragraph()
    p2.text = "• Point important 2"
    p2.level = 0
    p2_run = p2.runs[0]
    p2_run.font.size = Pt(24)
    p2_run.font.color.rgb = VERT_FONCE
    p2_run.font.name = 'Segoe UI'

    # Troisième paragraphe
    p3 = content_frame.add_paragraph()
    p3.text = "• Point important 3"
    p3.level = 0
    p3_run = p3.runs[0]
    p3_run.font.size = Pt(24)
    p3_run.font.color.rgb = VERT_FONCE
    p3_run.font.name = 'Segoe UI'

    # Footer avec logo/texte
    footer_shape2 = slide2.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(6.8), Inches(10), Inches(0.7)
    )
    footer_fill2 = footer_shape2.fill
    footer_fill2.solid()
    footer_fill2.fore_color.rgb = VERT_CLAIR
    footer_shape2.line.fill.background()

    footer_text2 = footer_shape2.text_frame
    footer_text2.text = "VOTRE VILLE, VOTRE VOIX • votrevillevotrevoix37.fr"
    footer_text2.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer_text2.vertical_anchor = MSO_ANCHOR.MIDDLE

    footer_run2 = footer_text2.paragraphs[0].runs[0]
    footer_run2.font.size = Pt(16)
    footer_run2.font.bold = True
    footer_run2.font.color.rgb = VERT_FONCE
    footer_run2.font.name = 'Segoe UI'


    # Slide 3: Deux colonnes
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # Fond blanc
    background3 = slide3.background
    fill3 = background3.fill
    fill3.solid()
    fill3.fore_color.rgb = BLANC

    # Bandeau supérieur
    header_shape3 = slide3.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    header_fill3 = header_shape3.fill
    header_fill3.solid()
    header_fill3.fore_color.rgb = VERT_FONCE
    header_shape3.line.fill.background()

    header_text3 = header_shape3.text_frame
    header_text3.text = "Titre sur deux colonnes"
    header_text3.paragraphs[0].alignment = PP_ALIGN.CENTER
    header_text3.vertical_anchor = MSO_ANCHOR.MIDDLE
    header_run3 = header_text3.paragraphs[0].runs[0]
    header_run3.font.size = Pt(36)
    header_run3.font.bold = True
    header_run3.font.color.rgb = BLANC
    header_run3.font.name = 'Segoe UI'

    # Colonne gauche avec fond vert clair
    left_col = slide3.shapes.add_shape(
        1, Inches(0.5), Inches(1.8), Inches(4.25), Inches(4.5)
    )
    left_fill = left_col.fill
    left_fill.solid()
    left_fill.fore_color.rgb = RGBColor(230, 245, 240)  # Vert très clair
    left_col.line.color.rgb = VERT_CLAIR
    left_col.line.width = Pt(2)

    left_text = left_col.text_frame
    left_text.word_wrap = True
    left_text.margin_left = Inches(0.3)
    left_text.margin_right = Inches(0.3)
    left_text.margin_top = Inches(0.3)

    left_p1 = left_text.paragraphs[0]
    left_p1.text = "Colonne de gauche"
    left_p1_run = left_p1.runs[0]
    left_p1_run.font.size = Pt(20)
    left_p1_run.font.bold = True
    left_p1_run.font.color.rgb = VERT_FONCE
    left_p1_run.font.name = 'Segoe UI'

    left_p2 = left_text.add_paragraph()
    left_p2.text = "\n• Élément 1\n• Élément 2\n• Élément 3"
    left_p2_run = left_p2.runs[0]
    left_p2_run.font.size = Pt(18)
    left_p2_run.font.color.rgb = VERT_FONCE
    left_p2_run.font.name = 'Segoe UI'

    # Colonne droite
    right_col = slide3.shapes.add_shape(
        1, Inches(5.25), Inches(1.8), Inches(4.25), Inches(4.5)
    )
    right_fill = right_col.fill
    right_fill.solid()
    right_fill.fore_color.rgb = RGBColor(230, 245, 240)
    right_col.line.color.rgb = VERT_CLAIR
    right_col.line.width = Pt(2)

    right_text = right_col.text_frame
    right_text.word_wrap = True
    right_text.margin_left = Inches(0.3)
    right_text.margin_right = Inches(0.3)
    right_text.margin_top = Inches(0.3)

    right_p1 = right_text.paragraphs[0]
    right_p1.text = "Colonne de droite"
    right_p1_run = right_p1.runs[0]
    right_p1_run.font.size = Pt(20)
    right_p1_run.font.bold = True
    right_p1_run.font.color.rgb = VERT_FONCE
    right_p1_run.font.name = 'Segoe UI'

    right_p2 = right_text.add_paragraph()
    right_p2.text = "\n• Élément A\n• Élément B\n• Élément C"
    right_p2_run = right_p2.runs[0]
    right_p2_run.font.size = Pt(18)
    right_p2_run.font.color.rgb = VERT_FONCE
    right_p2_run.font.name = 'Segoe UI'

    # Footer
    footer_shape3 = slide3.shapes.add_shape(
        1, Inches(0), Inches(6.8), Inches(10), Inches(0.7)
    )
    footer_fill3 = footer_shape3.fill
    footer_fill3.solid()
    footer_fill3.fore_color.rgb = VERT_CLAIR
    footer_shape3.line.fill.background()

    footer_text3 = footer_shape3.text_frame
    footer_text3.text = "VOTRE VILLE, VOTRE VOIX • votrevillevotrevoix37.fr"
    footer_text3.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer_text3.vertical_anchor = MSO_ANCHOR.MIDDLE
    footer_run3 = footer_text3.paragraphs[0].runs[0]
    footer_run3.font.size = Pt(16)
    footer_run3.font.bold = True
    footer_run3.font.color.rgb = VERT_FONCE
    footer_run3.font.name = 'Segoe UI'


    # Slide 4: Citation/Message important
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    # Fond dégradé (simulé avec vert moyen)
    background4 = slide4.background
    fill4 = background4.fill
    fill4.solid()
    fill4.fore_color.rgb = VERT_MOYEN

    # Citation centrale
    quote_box = slide4.shapes.add_textbox(
        Inches(1.5), Inches(2.5), Inches(7), Inches(2.5)
    )
    quote_frame = quote_box.text_frame
    quote_frame.word_wrap = True
    quote_frame.text = "« L'avenir, c'est vous »"
    quote_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    quote_run = quote_frame.paragraphs[0].runs[0]
    quote_run.font.size = Pt(48)
    quote_run.font.bold = True
    quote_run.font.italic = True
    quote_run.font.color.rgb = BLANC
    quote_run.font.name = 'Segoe UI'

    # Signature
    signature_box = slide4.shapes.add_textbox(
        Inches(1.5), Inches(5.2), Inches(7), Inches(0.8)
    )
    signature_frame = signature_box.text_frame
    signature_frame.text = "Brigitte Vengeon\nTête de liste VOTRE VILLE VOTRE VOIX"
    signature_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    for paragraph in signature_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(24)
            run.font.color.rgb = VERT_CLAIR
            run.font.name = 'Segoe UI'


    # Slide 5: Merci / Fin
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    # Fond vert foncé
    background5 = slide5.background
    fill5 = background5.fill
    fill5.solid()
    fill5.fore_color.rgb = VERT_FONCE

    # Titre MERCI
    merci_box = slide5.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(1.5)
    )
    merci_frame = merci_box.text_frame
    merci_frame.text = "MERCI"
    merci_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    merci_run = merci_frame.paragraphs[0].runs[0]
    merci_run.font.size = Pt(72)
    merci_run.font.bold = True
    merci_run.font.color.rgb = VERT_CLAIR
    merci_run.font.name = 'Segoe UI'

    # Informations de contact
    contact_box = slide5.shapes.add_textbox(
        Inches(1), Inches(4), Inches(8), Inches(2)
    )
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True

    p_site = contact_frame.paragraphs[0]
    p_site.text = "votrevillevotrevoix37.fr"
    p_site.alignment = PP_ALIGN.CENTER
    p_site_run = p_site.runs[0]
    p_site_run.font.size = Pt(28)
    p_site_run.font.bold = True
    p_site_run.font.color.rgb = BLANC
    p_site_run.font.name = 'Segoe UI'

    p_email = contact_frame.add_paragraph()
    p_email.text = "votrevillevotrevoix@gmail.com"
    p_email.alignment = PP_ALIGN.CENTER
    p_email_run = p_email.runs[0]
    p_email_run.font.size = Pt(22)
    p_email_run.font.color.rgb = BLANC
    p_email_run.font.name = 'Segoe UI'

    p_phone = contact_frame.add_paragraph()
    p_phone.text = "06 72 25 87 67"
    p_phone.alignment = PP_ALIGN.CENTER
    p_phone_run = p_phone.runs[0]
    p_phone_run.font.size = Pt(22)
    p_phone_run.font.color.rgb = BLANC
    p_phone_run.font.name = 'Segoe UI'

    p_vote = contact_frame.add_paragraph()
    p_vote.text = "\nMunicipales 15 et 22 mars 2026"
    p_vote.alignment = PP_ALIGN.CENTER
    p_vote_run = p_vote.runs[0]
    p_vote_run.font.size = Pt(20)
    p_vote_run.font.color.rgb = VERT_CLAIR
    p_vote_run.font.name = 'Segoe UI'

    return prs


if __name__ == '__main__':
    print("Création du template PowerPoint...")
    prs = create_presentation()
    output_file = 'Template-VVVV-Chateau-Renault-2026.pptx'
    prs.save(output_file)
    print(f"✅ Template créé avec succès : {output_file}")
    print("\nLe template contient 5 diapositives modèles :")
    print("  1. Page de titre")
    print("  2. Titre et contenu avec points")
    print("  3. Deux colonnes")
    print("  4. Citation/Message important")
    print("  5. Page de remerciement/fin")
    print("\nVous pouvez dupliquer et modifier ces diapositives selon vos besoins.")
